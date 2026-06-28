#!/usr/bin/env python3
"""Build the project summary slide deck (PowerPoint) mirroring
docs/bao_cao_tong_ket_phase3-6.md.

  python -m novbts.report.make_slides
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from novbts.paths import DOCS, RUNS, ensure

ACC = RGBColor(0x1E, 0x46, 0x90)
ACC2 = RGBColor(0x2E, 0x6D, 0xB4)
INK = RGBColor(0x18, 0x18, 0x18)
MUT = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xEE, 0xF2, 0xF8)

OUT = DOCS / "bao_cao_tong_ket_phase3-6.pptx"
SW, SH = Inches(13.333), Inches(7.5)


def _md_runs(para, text, size, color=INK, bold_all=False):
    """Add runs to a paragraph, honoring **bold** and `code` segments."""
    parts = re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", text)
    for p in parts:
        if not p:
            continue
        r = para.add_run()
        if p.startswith("**") and p.endswith("**"):
            r.text = p[2:-2]; r.font.bold = True; r.font.name = "Calibri"
        elif p.startswith("`") and p.endswith("`"):
            r.text = p[1:-1]; r.font.bold = bold_all; r.font.name = "Consolas"
        else:
            r.text = p; r.font.bold = bold_all; r.font.name = "Calibri"
        r.font.size = Pt(size); r.font.color.rgb = color


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_bar(slide, kicker, heading):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACC; bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(1.0)).text_frame
    tb.word_wrap = True
    if kicker:
        p = tb.paragraphs[0]; _md_runs(p, kicker, 13, ACC2, bold_all=True)
        p2 = tb.add_paragraph()
    else:
        p2 = tb.paragraphs[0]
    _md_runs(p2, heading, 26, ACC, bold_all=True)
    return slide


def body_box(slide, top=Inches(1.55), left=Inches(0.6), width=Inches(12.1), height=Inches(5.5)):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    return tf


def bullets(tf, items, size=16, gap=6):
    """items: list of (level, text)."""
    first = True
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        prefix = ("• " if level == 0 else "– ")
        _md_runs(p, prefix + text, size - level * 1)


def add_table(slide, headers, rows, left, top, width, col_w=None, fs=12, hfs=12.5):
    nr, nc = len(rows) + 1, len(headers)
    height = Inches(0.32 * nr)
    gtbl = slide.shapes.add_table(nr, nc, left, top, width, height)
    tbl = gtbl.table
    if col_w:
        tot = sum(col_w)
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Emu(int(width * w / tot))
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = ACC
        c.margin_top = Pt(2); c.margin_bottom = Pt(2)
        tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        _md_runs(p, h, hfs, WHITE, bold_all=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = BAND if i % 2 else WHITE
            c.margin_top = Pt(1); c.margin_bottom = Pt(1)
            tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
            _md_runs(p, str(val), fs)
    return gtbl


def note(slide, text, top=Inches(6.7), color=MUT, size=12):
    tf = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.6)).text_frame
    tf.word_wrap = True
    _md_runs(tf.paragraphs[0], text, size, color)


def pic(slide, path, left, top, width=None, height=None):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


# --------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    # 1 — title
    s = slide_blank(prs)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = ACC; bg.line.fill.background()
    bg.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(3.0)).text_frame
    tb.word_wrap = True
    _md_runs(tb.paragraphs[0], "Neural-Operator surrogate cho cảm biến xúc giác VBTS", 36, WHITE, bold_all=True)
    p = tb.add_paragraph(); _md_runs(p, "Báo cáo tổng kết dự án — Giai đoạn 3 → 6", 20, RGBColor(0xD8, 0xE2, 0xF2))
    p = tb.add_paragraph(); p.space_before = Pt(18)
    _md_runs(p, "FNO khả vi · cảm biến marker-dot · điều khiển · framework", 15, RGBColor(0xC4, 0xD4, 0xEC))
    p = tb.add_paragraph(); p.space_before = Pt(8)
    _md_runs(p, "2026-06-17", 14, RGBColor(0xC4, 0xD4, 0xEC))

    # 2 — executive summary
    s = slide_blank(prs); title_bar(s, "0 · TÓM TẮT ĐIỀU HÀNH", "Bốn kết quả cốt lõi (đều kiểm chứng trên FEM thật)")
    bullets(body_box(s), [
        (0, "**Framing field→field là chìa khóa** (P3): FNO thắng MLP per-point **2.24×** trên FEM thật (6.7× giải tích) — chuyển vị bề mặt là hàm **phi cục bộ**."),
        (0, "**Khả vi → điều khiển hiệu quả** (P4): học policy backprop qua FNO cần **~64× ít truy vấn / ~24× nhanh** hơn gradient-free (ES), cùng chất lượng."),
        (0, "**Cảm biến khả vi end-to-end** (P5): trường → ảnh marker-dot bằng torch; khôi phục lực kéo (sx,sy) **từ ảnh** sai số **2.1%**."),
        (0, "**Framework hoàn chỉnh** (P6): env khả vi (6a) + sàn nhiễu (6b) + lịch sử tải NULL (6c) + hình học vật (6d)."),
        (0, "**Trung thực:** trần tiếp tuyến ~0.15/14° là giới hạn nội tại; tốc độ không phải lợi thế bền; vs U-Net lợi thế hẹp còn ~1.03×."),
    ], size=16, gap=10)

    # 3 — context
    s = slide_blank(prs); title_bar(s, "1 · BỐI CẢNH & MỤC TIÊU", "Hai nhánh, ba lớp mô hình")
    bullets(body_box(s, height=Inches(2.2)), [
        (0, "Cảm biến VBTS: robot 'nhìn' tiếp xúc qua biến dạng gel in chấm marker. Mô phỏng chính xác = giải tiếp xúc đàn hồi-ma sát → đắt, không khả vi."),
        (0, "**Ý tưởng:** học một toán tử thay solver. Đích = ghép 2 nhánh:"),
        (1, "**Nhánh A** — môi trường mô phỏng cảm biến (GT + ảnh + env điều khiển)."),
        (1, "**Nhánh B** — FNO surrogate khả vi làm lõi nhanh trong vòng điều khiển."),
    ], size=15, gap=6)
    add_table(s, ["Lớp", "Vai trò", "Công cụ"], [
        ["GT chính xác-chậm", "sinh dữ liệu + thước đo", "PhysX FEM (Isaac Sim)"],
        ["Validator giải tích", "kiểm tra GT + baseline vật lý", "Hertz + Cattaneo–Mindlin"],
        ["Surrogate nhanh-khả-vi", "đóng góp chính", "FNO (Fourier Neural Operator)"],
    ], Inches(0.6), Inches(4.5), Inches(12.1), col_w=[3, 4, 4], fs=13)

    # 4 — foundation
    s = slide_blank(prs); title_bar(s, "2 · NỀN TẢNG", "Ground truth FEM & validator")
    bullets(body_box(s), [
        (0, "**Generator FEM** (`isaac_extract_{normal,shear}`, Docker `isaac-lab-fem`): gel 50×50×20 mm, indentor cầu r=0.02 m."),
        (0, "Xuất `disp[N,M,3]` (ux,uy,uz) bề mặt gel tại **lưới probe 32×32 = 1024 điểm** (lấy mẫu trường, độc lập mesh FEM)."),
        (0, "Shear/slip qua **micro-steps kéo tiếp tuyến** (phá deadlock tiếp xúc); ~3–4.6 s/frame (res-24)."),
        (0, "**Validator Hertz–Mindlin:** bán kính tiếp xúc lệch 1.3%, hội tụ 0.5% → GT FEM đáng tin ở chế độ pháp tuyến."),
        (0, "**4 chế độ tiếp xúc:** no-contact / stick / partial-slip / full-slip; quyết định không thêm mode."),
    ], size=16, gap=9)

    # 5 — P3 field2field
    s = slide_blank(prs); title_bar(s, "3 · PHASE 3 (GATE 3)", "Khám phá then chốt: framing field→field")
    bullets(body_box(s, height=Inches(3.0)), [
        (0, "**param→field** (mỗi điểm mớm full tham số): MLP per-point **thắng** FNO → bài toán thành cục bộ → chỉ là ablation phản chứng."),
        (0, "**field→field** (trường lún + shear·mask → trường dịch chuyển): MLP **sụp đổ** (0.743, hướng 63°), FNO giữ 0.111 → **FNO thắng 6.7×**."),
        (0, "**Luận điểm:** chuyển vị bề mặt phi cục bộ (Green's function) — FNO học tích chập phổ toàn cục, MLP cục bộ không thể."),
    ], size=16, gap=9)
    box = s.shapes.add_shape(1, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.4))
    box.fill.solid(); box.fill.fore_color.rgb = BAND; box.line.color.rgb = ACC2
    tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12); tf.margin_top = Pt(8)
    _md_runs(tf.paragraphs[0], "→ Đây là đóng góp khoa học chính của dự án: neural operator nắm cấu trúc phi-cục-bộ + phi-tuyến của trường slip mà mô hình vật lý/tuyến tính nhanh bỏ sót.", 16, ACC)

    # 6 — P3 Gate 3 results
    s = slide_blank(prs); title_bar(s, "3 · PHASE 3 (GATE 3)", "Gate 3 đóng trên FEM thật — 2400 frame")
    add_table(s, ["Chỉ số", "Kết quả"], [
        ["RQ1 — FNO overall relL2", "0.144–0.146 (hướng tiếp tuyến 14.6°)"],
        ["FNO thắng MLP per-point", "2.24× (0.146 vs 0.328; hướng 14.8° vs 35.7°)"],
        ["Slip-F1 head-a (multitask)", "0.904  (>0.75 ✓)"],
        ["Slip-F1 head-b (sau khi cứu)", "0.753  (+400 frame normal: F1 0→0.85)"],
        ["RQ2 — ngoại suy R/μ/E", "~1.3× cả ba trục (mượt, trong-hộp)"],
        ["RQ3 — tốc độ", "FNO ~7868 fps vs FEM 0.341 fps = ~23.000×"],
    ], Inches(0.6), Inches(1.7), Inches(12.1), col_w=[5, 7], fs=14)
    note(s, "Luận điểm phi-cục-bộ giữ vững trên vật lý thật (giải tích 6.7× → FEM 2.24×, hẹp hơn vì FEM nhiễu/ít lý tưởng).", top=Inches(5.0))

    # 7 — P3 bake-off
    s = slide_blank(prs); title_bar(s, "3 · PHASE 3 — BAKE-OFF", "So với các mô phỏng VBTS tiêu biểu (cô lập giá trị)")
    add_table(s, ["Phương pháp", "overall relL2", "hướng°", "FNO hơn"], [
        ["TACTO-style (động học, no friction)", "0.504", "65.8", "3.51×"],
        ["Cattaneo–Mindlin giải tích (đã hiệu chỉnh)", "0.435", "39.0", "3.03×"],
        ["Taxim/FOTS-style (tuyến tính chồng chập)", "0.295", "26.2", "2.05×"],
        ["MLP per-point (cục bộ)", "0.321", "33.6", "2.24×"],
        ["**FNO (ours)**", "**0.144**", "**14.6**", "—"],
    ], Inches(0.6), Inches(1.7), Inches(12.1), col_w=[6, 2.5, 1.8, 1.7], fs=13)
    bullets(body_box(s, top=Inches(4.6), height=Inches(2.3)), [
        (0, "Vật lý giải tích kinh điển (Mindlin) **thua** mô hình tuyến tính fit-data — affine không đổi được hình dạng profile."),
        (0, "FNO hơn vì chuyển stick→partial→full slip là **phi tuyến**; MLP cục bộ sụp ở slip."),
        (0, "**Trung thực:** vs mạng SOTA (U-Net 0.148 ≈ ngang, ít hơn 5.7× params) lợi thế thu hẹp còn 1.03–1.46×."),
    ], size=14, gap=6)

    # 8 — P4
    s = slide_blank(prs); title_bar(s, "4 · PHASE 4", "Điều khiển qua FNO khả vi (autograd vs ES)")
    add_table(s, ["Phương pháp", "final loss", "truy vấn (fwd)", "wall"], [
        ["autograd (backprop qua FNO)", "8.02e-8", "300", "13 s"],
        ["ES (gradient-free, pop 32)", "8.36e-8", "19 200", "313 s"],
    ], Inches(0.6), Inches(1.7), Inches(8.0), col_w=[5, 3, 3, 2], fs=13)
    pic(s, RUNS / "phase4" / "policy_servo_curve.png", Inches(8.9), Inches(1.7), width=Inches(3.9))
    bullets(body_box(s, top=Inches(3.4), width=Inches(8.0), height=Inches(3.6)), [
        (0, "Cùng chất lượng cuối, autograd **~64× ít truy vấn / ~24× nhanh hơn**."),
        (0, "FNO **đóng băng**, gradient chỉ chảy tới action; là policy ngữ-cảnh một-bước (amortized)."),
        (0, "**Phản đề Suh 2022 không cắn:** FNO làm trơn bước nhảy stick→slip → gradient sạch (dở cho forward, tốt cho gradient)."),
        (0, "Task A (anti-slip) bỏ trung thực: proxy nhiễu (tăng theo depth = hiệu ứng diện tích)."),
    ], size=14, gap=6)

    # 9 — P5 sensor
    s = slide_blank(prs); title_bar(s, "5 · PHASE 5", "Mô hình cảm biến marker-dot khả vi")
    bullets(body_box(s, height=Inches(2.5)), [
        (0, "Ảnh cảm biến = **chiếu + render chấm** của trường có sẵn (không giải FEM mới); torch → `render∘FNO` khả vi."),
        (0, "**Pipeline:** disp[3,32,32] → lấy mẫu 121 chấm → chiếu camera (pinhole dưới màng) → render Gaussian-splat → ảnh."),
        (0, "ux,uy → chấm trượt ngang (marker flow); uz → đổi depth → chấm phóng to/nhỏ (độ sâu lún)."),
    ], size=15, gap=7)
    add_table(s, ["Chỉ số", "Giá trị"], [
        ["Faithful: cos(flow, disp_xy)", "0.973"],
        ["Compat FNO+renderer (marker-flow)", "0.26 rel-L2 (khớp trần tiếp tuyến FNO)"],
        ["Inverse từ ẢNH (sx,sy qua render∘FNO)", "2.1% / 1.1° → pipeline khả vi end-to-end chạy"],
    ], Inches(0.6), Inches(4.4), Inches(12.1), col_w=[5, 7], fs=13)

    # 10 — P6 overview
    s = slide_blank(prs); title_bar(s, "6 · PHASE 6", "Hoàn thiện framework — 4 hướng")
    add_table(s, ["Hướng", "Công việc", "Kết quả cốt lõi"], [
        ["6a env wrapper [TRỌNG TÂM]", "gói FNO+sensor+reward thành env khả vi", "policy đóng 87% gap, gradcheck OK"],
        ["6b realism/sim2real", "nhiễu camera + tracker → đo sàn nhiễu", "FNO 0.54px > nhiễu 0.06px"],
        ["6c temporal + lịch sử tải", "render quỹ đạo + test path-dependence", "trạng thái cuối chỉ phụ thuộc endpoint (NULL)"],
        ["6d hình học vật", "indentor sphere/flat/cylinder (Isaac)", "flow 1.7 / 10.5 / 8.4 px khác hẳn"],
    ], Inches(0.6), Inches(1.8), Inches(12.1), col_w=[3, 4.5, 4.5], fs=13)

    # 11 — P6a env (with image)
    s = slide_blank(prs); title_bar(s, "6a · ENV KHẢ VI [TRỌNG TÂM]", "Một API gói cả forward + sensor + control")
    bullets(body_box(s, width=Inches(6.6), height=Inches(4.0)), [
        (0, "`reset()` → context sphere + render **target imprint** (mục tiêu)."),
        (0, "`differentiable_step(sx,sy)` → FNO → render → ảnh; **reward = −MSE(ảnh, target)**, khả vi."),
        (0, "Tích hợp **PolicyMLP P4**, huấn luyện qua env. Adapter gymnasium."),
        (0, "Single-step contextual (FNO tĩnh); chỉ dùng sphere."),
    ], size=14, gap=6)
    add_table(s, ["", "mean reward (cao=tốt)"], [
        ["ngẫu nhiên", "−3.23e-3"],
        ["Phase-4 policy", "−6.52e-4"],
        ["oracle (thật)", "−2.72e-4"],
    ], Inches(0.6), Inches(5.2), Inches(6.4), col_w=[3, 4], fs=13)
    pic(s, RUNS / "phase6" / "env_demo.png", Inches(7.5), Inches(1.7), height=Inches(5.3))
    note(s, "Policy đóng 87% gap ngẫu-nhiên→oracle · gradcheck rel_err 3.7% · cột policy khớp target →", top=Inches(7.05), size=11)

    # 12 — 6b realism
    s = slide_blank(prs); title_bar(s, "6b · REALISM / SIM2REAL", "Sàn nhiễu camera vs sai số FNO")
    bullets(body_box(s, width=Inches(7.2), height=Inches(4.5)), [
        (0, "Thang đo chuẩn = **EPE pixel**. GT-vs-FNO: EPE **0.54 px** (pitch marker 12 px)."),
        (0, "Quét read-noise 0→8%: jitter (sàn nhiễu) 0.055→0.107 px; bias (tracker) 0.31 px."),
        (0, "**Kết quả (đảo kỳ vọng):** sai số FNO 0.54 px **> sàn nhiễu @2% (0.06)** và **> tracker (0.31)**."),
        (0, "→ Nhiễu camera **không che được** sai số FNO. Muốn cải thiện: nâng FNO / tracker sub-pixel, đừng đổ tiền camera."),
        (0, "calibration.py: chỉ schema (chờ hardware thật)."),
    ], size=14, gap=7)
    pic(s, RUNS / "phase6" / "realism_floor.png", Inches(8.1), Inches(1.8), width=Inches(4.7))

    # 13 — 6c temporal
    s = slide_blank(prs); title_bar(s, "6c · TEMPORAL + LỊCH SỬ TẢI", "Trạng thái cuối có phụ thuộc đường tải?")
    bullets(body_box(s, width=Inches(7.2), height=Inches(4.5)), [
        (0, "Generator thêm `--save-trajectory` + 3 đường tải (linear/ortho/reverse) cùng endpoint."),
        (0, "Render quỹ đạo kéo → **video marker-dot** (progressive stick→slip); reverse hiện overshoot ở f≈0.71."),
        (0, "**Test path-dependence (NULL):** endpoint-only 0.161 vs +load-mode 0.171; kNN distance-matched cross/same = **1.00**."),
        (0, "→ Trạng thái cuối **chỉ phụ thuộc endpoint (sx,sy,depth)**, không phụ thuộc đường đi."),
        (0, "→ Thành phần bất khả giảm của trần tiếp tuyến **không** phải lịch sử tải."),
    ], size=14, gap=7)
    pic(s, RUNS / "phase6" / "temporal_slip_curve.png", Inches(8.1), Inches(1.9), width=Inches(4.7))

    # 14 — 6d geometry
    s = slide_blank(prs); title_bar(s, "6d · HÌNH HỌC VẬT", "Mở rộng beyond sphere (Isaac)")
    add_table(s, ["Hình", "mean_tang (m)", "contact-area", "marker flow"], [
        ["sphere", "0.0004", "39%", "1.7 px"],
        ["cylinder", "0.0016", "93%", "8.4 px"],
        ["flat punch", "0.0022", "100%", "10.5 px"],
    ], Inches(0.6), Inches(1.7), Inches(7.0), col_w=[3, 3, 3, 3], fs=13)
    bullets(body_box(s, top=Inches(3.6), width=Inches(7.0), height=Inches(3.2)), [
        (0, "**Hình học định hình mạnh trường xúc giác** (flow khác ~6×)."),
        (0, "flat punch phủ gần hết gel; cylinder chậm (~30 s/frame)."),
        (0, "mesh (UsdFileCfg) đã wire nhưng chưa test (cần file USD)."),
    ], size=14, gap=6)
    pic(s, RUNS / "phase6" / "object_geometry_slip.png", Inches(8.0), Inches(1.9), width=Inches(4.8))

    # 15 — tangential ceiling
    s = slide_blank(prs); title_bar(s, "7 · PHÁT HIỆN XUYÊN SUỐT", "Trần tiếp tuyến ~0.15 / 14° — 8 đòn bẩy đều âm tính")
    add_table(s, ["Đòn bẩy", "Kết quả"], [
        ["Data-scaling 200→1600 frame", "bão hòa ~1600 (gấp đôi chỉ mua ~0.01)"],
        ["Model capacity (3× params)", "0.146→0.151 (không giúp)"],
        ["Lưới mịn (res-24→res-32)", "0.146→0.158 (không hạ)"],
        ["Nhiều Fourier modes (12→16→20)", "tệ hơn / lỗi (>Nyquist)"],
        ["Input: mớm trường Mindlin", "NULL (hàm tất định của input)"],
        ["Input: ranh giới sắc / U-Net", "lợi nhỏ ~5% (vá phần FNO làm trơn)"],
        ["Lịch sử tải (6c)", "NULL (chỉ phụ thuộc endpoint)"],
    ], Inches(0.6), Inches(1.7), Inches(12.1), col_w=[6, 6], fs=12.5)
    note(s, "Trần = (1) sắc-nhưng-FNO-trơn (~5%, vá nhỏ) + (2) bất khả giảm từ tham số macro (phần lớn). KHÔNG do GT/mesh/model. Pháp tuyến đã hội tụ tốt.", top=Inches(6.5), size=12)

    # 16 — limitations
    s = slide_blank(prs); title_bar(s, "8 · HẠN CHẾ & TRUNG THỰC", "Ghi rõ để không overclaim")
    bullets(body_box(s), [
        (0, "**Env single-step:** FNO là map tĩnh → chưa mô phỏng động lực nhiều bước."),
        (0, "**Tốc độ không phải lợi thế bền:** sim vật lý GPU (Taccel) throughput vượt FNO; lợi thế = latency + khả-vi + nhẹ-HW + mesh-free."),
        (0, "**Lợi thế vs SOTA hẹp:** U-Net ≈ ngang FNO."),
        (0, "**Trần tiếp tuyến** không phá được bằng các đòn bẩy đã thử."),
        (0, "**Chưa có hardware:** calibration mới là schema; sim2real chưa đóng vòng."),
        (0, "Hình học mesh chưa test; một số kết quả single-seed cần xác nhận."),
    ], size=16, gap=9)

    # 17 — framework + roadmap
    s = slide_blank(prs); title_bar(s, "9 · FRAMEWORK & LỘ TRÌNH", "Trạng thái hiện tại + bước tiếp")
    box = s.shapes.add_shape(1, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.0))
    box.fill.solid(); box.fill.fore_color.rgb = BAND; box.line.color.rgb = ACC2
    tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Pt(14); tf.margin_top = Pt(10)
    _md_runs(tf.paragraphs[0], "FRAMEWORK HIỆN CÓ (chạy ngay, pure-Python + GPU)", 15, ACC, bold_all=True)
    for t in ["reset() / step() / differentiable_step()  ← tactile_env.py (6a)",
              "  ├─ FNO surrogate (đóng băng) · marker-dot sensor · camera noise + tracker · PolicyMLP"]:
        p = tf.add_paragraph(); _md_runs(p, t, 14, INK)
    bullets(body_box(s, top=Inches(4.0), height=Inches(3.0)), [
        (0, "**1. Env Isaac-Sim thật** với robot gắn cảm biến — FNO làm lõi nhanh trong vòng điều khiển (mảnh lớn còn lại Nhánh A)."),
        (0, "**2. Multi-step dynamics** — nâng env nhiều bước (FNO có thời gian)."),
        (0, "**3. Sim2real** — dựng cảm biến DIY, fit SensorCalib, fine-tune FNO trên dữ liệu thật."),
        (0, "**4. Hình học mesh** vật thật (file USD)."),
    ], size=15, gap=8)

    # 18 — closing
    s = slide_blank(prs)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = ACC; bg.line.fill.background(); bg.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.5)).text_frame
    tb.word_wrap = True
    _md_runs(tb.paragraphs[0], "Một mô phỏng VBTS khả vi, dùng được", 32, WHITE, bold_all=True)
    p = tb.add_paragraph(); p.space_before = Pt(14)
    _md_runs(p, "Fidelity-cỡ-FEM · khả vi end-to-end · nhẹ phần cứng · mesh-free — sẵn sàng cho điều khiển/RL và sim2real.", 17, RGBColor(0xD8, 0xE2, 0xF2))
    p = tb.add_paragraph(); p.space_before = Pt(18)
    _md_runs(p, "Chi tiết: docs/bao_cao_tong_ket_phase3-6.md (+ PDF)", 14, RGBColor(0xC4, 0xD4, 0xEC))

    ensure(DOCS)
    prs.save(str(OUT))
    print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
